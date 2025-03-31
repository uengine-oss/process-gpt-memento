from fastapi import FastAPI, Request, Body
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

import requests
import tempfile
import fitz  # PyMuPDF
import docx
import pptx
import os

from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import SupabaseVectorStore
from supabase import create_client

from langchain.chains import RetrievalQA
from langchain.chat_models import ChatOpenAI

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # 모든 출처 허용
    allow_credentials=True,
    allow_methods=["*"],  # 모든 HTTP 메서드 허용
    allow_headers=["*"],  # 모든 HTTP 헤더 허용
)

supabase_url = os.getenv("SUPABASE_URL")
if not supabase_url:
    supabase_url = "http://127.0.0.1:54321"

supabase_key = os.getenv("SUPABASE_KEY")
if not supabase_key:
    supabase_key = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZS1kZW1vIiwicm9sZSI6InNlcnZpY2Vfcm9sZSIsImV4cCI6MTk4MzgxMjk5Nn0.EGIM96RAZx35lJzdJsyH-qQwv8Hdp7fsn3W0YpN81IU"

supabase_client = create_client(supabase_url, supabase_key)

openai_api_key = os.getenv("OPENAI_API_KEY")
model = ChatOpenAI(openai_api_key=openai_api_key)
embeddings = OpenAIEmbeddings(model="text-embedding-3-small", deployment="text-embedding-3-small", openai_api_key=openai_api_key)

vector_store = SupabaseVectorStore(
    client=supabase_client,
    embedding=embeddings,
    table_name="documents",
    query_name="match_documents",
)


class FileRequest(BaseModel):
    path: str

@app.post("/index")
async def index_file(request: Request):
    json_data = await request.json()
    path = json_data.get('path')
    
    if path is None:
        return {"message": "Path is required."}
    
    host_name = request.headers.get('X-Forwarded-Host')
    if host_name is None or any(substring in host_name for substring in ['localhost']):
        tenant_id = 'localhost'
    else:
        tenant_id = host_name.split('.')[0]
    
    file_name = path.split("/")[-1]
    file_url = f"{supabase_url}/storage/v1/object/public/{path}"
    file_ext = path.split(".")[-1]

    response = requests.get(file_url)
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
        tmp.write(response.content)
        tmp_path = tmp.name

    content = extract_text(tmp_path, file_ext)
    if not content.strip():
        return {"message": "No text extracted from file."}

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100
    )
    chunks = text_splitter.split_text(content)

    documents = [
        Document(
            page_content=chunk,
            metadata={
                "tenant_id": tenant_id,
                "type": "upload_file",
                "file_path": file_url,
                "file_name": file_name,
                "chunk_index": i
            }
        )
        for i, chunk in enumerate(chunks)
    ]

    vector_store.add_documents(documents)

    return {"message": f"{len(chunks)} chunks indexed."}

def extract_text(path: str, ext: str) -> str:
    if ext == "pdf":
        return "\n".join([page.get_text() for page in fitz.open(path)])
    elif ext == "docx":
        doc = docx.Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    elif ext == "pptx":
        ppt = pptx.Presentation(path)
        return "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif ext == "md":
        return open(path, "r", encoding="utf-8").read()
    else:
        return {"message": f"Unsupported file format: {ext}"}



@app.post("/query")
def query_vector_store(input: dict = Body(...)):
    query = input.get("query", "")
    tenant_id = input.get("tenant_id", "localhost")
    
    retriever = vector_store.as_retriever(search_kwargs={
        "k": 5,
        "filter": {
            "tenant_id": tenant_id,
            "type": "upload_file"
        }
    })
    
    qa_chain = RetrievalQA.from_chain_type(
        llm=model,
        retriever=retriever,
        return_source_documents=True
    )

    result = qa_chain.invoke(query)

    return {
        "response": result["result"],
        "metadata": {
            f"{doc.metadata.get('file_name', 'unknown')}#{doc.metadata.get('chunk_index', i)}": doc.metadata
            for i, doc in enumerate(result["source_documents"])
        }
    }


import uvicorn

if __name__ == "__main__":
    uvicorn.run("memento-service:app", host="0.0.0.0", port=8005, log_level="info")

